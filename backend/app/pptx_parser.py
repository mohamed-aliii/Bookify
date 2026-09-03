import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class SlideTextBlock:
    text: str
    kind: str  # title | subtitle | body | note


@dataclass
class SlideContent:
    slide_number: int
    title: str | None
    bullets: list[str] = field(default_factory=list)
    body_text: str = ""
    notes: str | None = None
    raw_text: str = ""


@dataclass
class ParsedSlides:
    title: str
    num_pages: int
    slides: list[SlideContent]


def _shape_text(shape) -> str | None:
    if not shape.has_text_frame:
        return None
    text = shape.text.strip()
    return text if text else None


def _extract_slide_title(slide) -> str | None:
    # Try title placeholder first
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        t = slide.shapes.title.text.strip()
        if t:
            return t
    # Fallback: first shape with placeholder type TITLE
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.type.name == "TITLE" and shape.has_text_frame:
                t = shape.text.strip()
                if t:
                    return t
        except Exception:
            continue
    # Fallback: largest text on slide as title (heuristic)
    candidates: list[tuple[int, str]] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if not t:
            continue
        # skip very long bodies as title candidate
        if len(t) > 200:
            continue
        # estimate size from first paragraph font size if available
        size = 0
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        size = max(size, run.font.size.pt or 0)
        except Exception:
            pass
        candidates.append((size, t))
    if candidates:
        candidates.sort(reverse=True)
        # If largest is at least 14pt or clearly larger, use it
        if candidates[0][0] >= 14 or len(candidates[0][1].split()) <= 12:
            return candidates[0][1]
    return None


def _extract_bullets_and_body(slide, title_text: str | None) -> tuple[list[str], str]:
    bullets: list[str] = []
    bodies: list[str] = []
    for shape in slide.shapes:
        if shape.has_table:
            # Extract table as text
            try:
                tbl = shape.table
                rows = []
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    cells = [c for c in cells if c]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    bodies.append("\n".join(rows))
            except Exception:
                continue
            continue
        if not shape.has_text_frame:
            continue
        # Skip the title itself
        if title_text and shape.text.strip() == title_text:
            continue
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if not txt:
                continue
            # Skip duplicate of title
            if title_text and txt == title_text:
                continue
            is_bullet = bool(para.level is not None and para.level > 0) or txt.startswith(("•", "-", "–", "·", "*"))
            # Also treat short lines as bullets (heuristic for slides)
            if is_bullet or (len(txt) < 220 and txt and not txt.endswith(".")):
                # Clean bullet marker
                clean = txt.lstrip("•-–·* ").strip()
                if clean:
                    bullets.append(clean)
            else:
                bodies.append(txt)
            # Also capture sub-bullets as body/bullets already
    body_text = "\n".join(bodies).strip()
    return bullets, body_text


def _extract_notes(slide) -> str | None:
    try:
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            t = notes_slide.notes_text_frame.text.strip()
            return t if t else None
    except Exception:
        pass
    return None


def parse_pptx(path: Path | str, fallback_title: str | None = None) -> ParsedSlides:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx is required for PPTX support. Install with: pip install python-pptx") from e

    path = Path(path)
    prs = Presentation(str(path))

    # Book title from presentation core properties or filename
    title = None
    try:
        if prs.core_properties.title:
            t = prs.core_properties.title.strip()
            if t and len(t) >= 3:
                title = t
    except Exception:
        pass
    if not title and fallback_title:
        import re

        fb = re.sub(r"[_\-.]+", " ", fallback_title).strip()
        fb = re.sub(r"\s+", " ", fb).strip()
        if fb:
            title = fb
    if not title:
        title = path.stem

    slides: list[SlideContent] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title_text = _extract_slide_title(slide)
        bullets, body_text = _extract_bullets_and_body(slide, title_text)
        notes = _extract_notes(slide)

        # Combine for raw_text (for chunking/AI)
        parts: list[str] = []
        if title_text:
            parts.append(title_text)
        if bullets:
            parts.extend([f"• {b}" for b in bullets])
        if body_text:
            parts.append(body_text)
        if notes:
            parts.append(f"[Notes] {notes}")
        raw = "\n".join(parts).strip()
        # Fallback: if extraction failed, grab any text on slide
        if not raw:
            fallback_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    if t:
                        fallback_parts.append(t)
            raw = "\n".join(fallback_parts).strip()
        # Ensure title fallback for unnamed slides
        effective_title = title_text or f"Slide {idx}"
        if not raw:
            raw = effective_title

        slides.append(
            SlideContent(
                slide_number=idx,
                title=effective_title,
                bullets=bullets,
                body_text=body_text,
                notes=notes,
                raw_text=raw,
            )
        )

    return ParsedSlides(title=title, num_pages=len(slides), slides=slides)
