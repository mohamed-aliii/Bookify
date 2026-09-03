import logging
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)


def is_libreoffice_available() -> bool:
    return shutil.which("soffice") is not None or shutil.which("libreoffice") is not None


def _libreoffice_bin() -> str | None:
    for b in ("soffice", "libreoffice"):
        p = shutil.which(b)
        if p:
            return p
    return None


def _convert_via_powerpoint(pptx_path: Path, dest_pdf: Path, timeout: int = 90) -> bool:
    """Convert PPTX to PDF using native Microsoft PowerPoint COM on Windows."""
    if sys.platform != "win32":
        return False
    ps_script = f"""
$ErrorActionPreference = 'Stop'
$ppt = New-Object -ComObject PowerPoint.Application
try {{
    $pres = $ppt.Presentations.Open('{str(pptx_path.resolve())}', 0, 0, 0)
    $pres.SaveAs('{str(dest_pdf.resolve())}', 32)
    $pres.Close()
}} finally {{
    $ppt.Quit()
    [System.Runtime.InteropServices.Marshal]::ReleaseComObject($ppt) | Out-Null
}}
"""
    with tempfile.NamedTemporaryFile(suffix=".ps1", delete=False, mode="w", encoding="utf-8") as f:
        f.write(ps_script)
        script_path = Path(f.name)
    try:
        cmd = [
            "powershell",
            "-NoProfile",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(script_path),
        ]
        res = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        if res.returncode == 0 and dest_pdf.exists() and dest_pdf.stat().st_size > 0:
            logger.info("Successfully converted PPTX to PDF via PowerPoint COM: %s", dest_pdf)
            return True
        logger.warning(
            "PowerPoint COM conversion failed: returncode=%s stdout=%s stderr=%s",
            res.returncode,
            res.stdout.strip(),
            res.stderr.strip(),
        )
        return False
    except Exception as e:
        logger.warning("PowerPoint COM conversion error: %s", e)
        return False
    finally:
        if script_path.exists():
            script_path.unlink()


def _convert_via_libreoffice(pptx_path: Path, expected_pdf: Path, timeout: int = 120) -> bool:
    """Convert PPTX to PDF using LibreOffice headless."""
    bin_path = _libreoffice_bin()
    if not bin_path:
        return False

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        cmd = [
            bin_path,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmp_path),
            str(pptx_path),
        ]
        logger.info("Converting PPTX to PDF via LibreOffice: %s", " ".join(cmd))
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
            if result.returncode != 0:
                logger.warning("LibreOffice conversion failed (%s): %s %s", result.returncode, result.stderr, result.stdout)
                return False
            tmp_pdf = tmp_path / f"{pptx_path.stem}.pdf"
            if not tmp_pdf.exists():
                candidates = list(tmp_path.glob("*.pdf"))
                if not candidates:
                    return False
                tmp_pdf = candidates[0]

            if expected_pdf.exists():
                expected_pdf.unlink()
            shutil.move(str(tmp_pdf), str(expected_pdf))
            return expected_pdf.exists() and expected_pdf.stat().st_size > 0
        except Exception as e:
            logger.warning("LibreOffice conversion error: %s", e)
            return False


def _convert_via_pure_python(pptx_path: Path, dest_pdf: Path) -> bool:
    """Fallback: convert PPTX slides to PDF using python-pptx and pymupdf."""
    try:
        import pymupdf as fitz
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(pptx_path))
        doc = fitz.open()
        sw = float(prs.slide_width.pt) if hasattr(prs.slide_width, "pt") else 960.0
        sh = float(prs.slide_height.pt) if hasattr(prs.slide_height, "pt") else 540.0

        for slide in prs.slides:
            page = doc.new_page(width=sw, height=sh)
            page.draw_rect(fitz.Rect(0, 0, sw, sh), color=None, fill=(0.98, 0.98, 0.99))

            for shape in slide.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_bytes = shape.image.blob
                        x0 = float(shape.left.pt)
                        y0 = float(shape.top.pt)
                        w = float(shape.width.pt)
                        h = float(shape.height.pt)
                        page.insert_image(fitz.Rect(x0, y0, x0 + w, y0 + h), stream=img_bytes)
                    elif shape.has_text_frame:
                        txt = shape.text_frame.text.strip()
                        if txt:
                            x0 = float(shape.left.pt)
                            y0 = float(shape.top.pt)
                            w = float(shape.width.pt)
                            h = float(shape.height.pt)
                            is_title = shape == slide.shapes.title
                            fs = 20 if is_title else 12
                            col = (0.1, 0.1, 0.15) if is_title else (0.2, 0.22, 0.25)
                            page.insert_textbox(fitz.Rect(x0, y0, x0 + w, y0 + h), txt, fontsize=fs, color=col)
                except Exception:
                    pass

        doc.save(str(dest_pdf))
        doc.close()
        return dest_pdf.exists() and dest_pdf.stat().st_size > 0
    except Exception as e:
        logger.warning("Pure-python PPTX conversion error: %s", e)
        return False


def convert_pptx_to_pdf(pptx_path: Path | str, dest_dir: Path | None = None) -> Path:
    """Convert PPTX to PDF via PowerPoint COM, LibreOffice, or pure-python fallback.

    Returns path to the generated PDF.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    dest_dir = Path(dest_dir) if dest_dir else pptx_path.parent
    dest_dir.mkdir(parents=True, exist_ok=True)
    expected_pdf = dest_dir / f"{pptx_path.stem}.pdf"

    # Tier 1: Microsoft PowerPoint on Windows
    if sys.platform == "win32":
        if _convert_via_powerpoint(pptx_path, expected_pdf):
            return expected_pdf

    # Tier 2: LibreOffice headless
    if is_libreoffice_available():
        if _convert_via_libreoffice(pptx_path, expected_pdf):
            return expected_pdf

    # Tier 3: Pure python fallback
    if _convert_via_pure_python(pptx_path, expected_pdf):
        return expected_pdf

    raise RuntimeError(
        f"Failed to convert {pptx_path.name} to PDF preview via PowerPoint, LibreOffice, or pure-python fallback."
    )
